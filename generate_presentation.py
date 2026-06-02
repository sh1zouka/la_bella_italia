#!/usr/bin/env python3
"""
Generate presentation for La Bella Italia diploma project.
12 slides, widescreen 16:9, lots of graphics, minimal text.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

RED = RGBColor(0xCC, 0x29, 0x2B)
GOLD = RGBColor(0xD4, 0x9B, 0x2C)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
CARD1 = RGBColor(0xFA, 0xFA, 0xFA)
BORDER = RGBColor(0xE8, 0xE8, 0xE8)

SW = Emu(12192000)  # slide width
SH = Emu(6858000)   # slide height

def bg(slide, color):
    b = slide.background.fill
    b.solid()
    b.fore_color.rgb = color

def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    return bx

def lines(slide, l, t, w, h, data):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, d in enumerate(data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = d.get('t', '')
        p.font.size = Pt(d.get('s', 16))
        p.font.color.rgb = d.get('c', WHITE)
        p.font.bold = d.get('b', False)
        p.space_after = Pt(d.get('sp', 4))
    return bx

def placeholder(slide, l, t, w, h, label):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    s.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    s.line.width = Pt(1.5)
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = LGRAY
    p.alignment = PP_ALIGN.CENTER
    return s

# ===== CREATE =====
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

# SLIDE 1 - Title
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK)
rect(s, 0, 0, SW, Emu(80000), RED)
rect(s, 0, 0, Emu(200000), SH, RED)
tb(s, Emu(500000), Emu(400000), Emu(3000000), Emu(1200000), "Pizza", sz=72)
tb(s, Emu(500000), Emu(1600000), Emu(11000000), Emu(1600000),
   "Razrabotka veb-prilozhenija dlja avtomatizacii\nprioma i obrabotki zakazov\nv italjanskom restorane \"La Bella Italia\"",
   sz=32, b=True)
rect(s, Emu(500000), Emu(3400000), Emu(3000000), Emu(40000), GOLD)
lines(s, Emu(500000), Emu(3600000), Emu(8000000), Emu(2400000), [
    {'t': 'Student: Kostja Suetov', 's': 20, 'c': RGBColor(0xCC,0xCC,0xCC)},
    {'t': 'Group: 9-421', 's': 18, 'c': RGBColor(0x99,0x99,0x99)},
    {'t': 'Specialty: 09.02.07 IS and Programming', 's': 16, 'c': RGBColor(0x99,0x99,0x99)},
    {'t': 'Supervisor: [Name]', 's': 18, 'c': RGBColor(0xCC,0xCC,0xCC)},
    {'t': '2026', 's': 22, 'c': GOLD, 'b': True},
])

# SLIDE 2 - Relevance
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Aktualnost", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
lines(s, Emu(500000), Emu(1400000), Emu(5500000), Emu(4800000), [
    {'t': 'Rynok dostavy edy v Rossii rastjot na 25% ezhegodno', 's': 18, 'c': DARK, 'b': True},
    {'t': '', 's': 10},
    {'t': '70% restoranov ispolzujut cifrovye kanaly zakaza', 's': 16, 'c': DARK},
    {'t': '', 's': 8},
    {'t': 'Spros na udobnye veb-prilozhenija postojanno rastjot', 's': 16, 'c': DARK},
    {'t': '', 's': 8},
    {'t': 'Nebolshie restorany nuzhdajutsja v dostupnyh reshenijah', 's': 16, 'c': DARK},
    {'t': '', 's': 8},
    {'t': 'Bonusnye sistemy povyshajut lojalnost na 30%', 's': 16, 'c': DARK},
])
tb(s, Emu(6500000), Emu(1500000), Emu(5000000), Emu(1000000), "25%", sz=60, b=True, c=RED, a=PP_ALIGN.CENTER)
tb(s, Emu(6500000), Emu(2400000), Emu(5000000), Emu(500000), "rost rynka\nezhegodno", sz=14, c=LGRAY, a=PP_ALIGN.CENTER)
tb(s, Emu(6500000), Emu(3200000), Emu(5000000), Emu(1000000), "70%", sz=60, b=True, c=GOLD, a=PP_ALIGN.CENTER)
tb(s, Emu(6500000), Emu(4100000), Emu(5000000), Emu(500000), "restoranov s\ncifrovymi kanalami", sz=14, c=LGRAY, a=PP_ALIGN.CENTER)

# SLIDE 3 - Goal
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Cel raboty", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
card = rect(s, Emu(500000), Emu(1500000), Emu(11000000), Emu(2000000), RGBColor(0xF5,0xF0,0xE0))
tb(s, Emu(700000), Emu(1700000), Emu(10600000), Emu(1600000),
   "Razrabotka veb-prilozhenija dlja avtomatizacii prijoma i obrabotki\nzakazov v italjanskom restorane \"La Bella Italia\"\ns ispolzovaniem sovremennyh veb-tehnologij",
   sz=20, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, Emu(500000), Emu(3800000), Emu(11000000), Emu(600000), "Arhitektura prilozhenija", sz=24, b=True, c=DARK)
cols = [
    ("Frontend", "HTML + CSS + JS\nlocalStorage\nAdaptivnyj dizajn", RGBColor(0xCC,0x29,0x2B)),
    ("Backend", "Node.js + Express\nNeDB-obortka\nREST API", RGBColor(0x2E,0x7D,0x32)),
    ("Database", "PostgreSQL 16\nJSONB-polja\nSequelize ORM", RGBColor(0x1A,0x5C,0x8A)),
]
for i, (title, desc, accent) in enumerate(cols):
    x = Emu(500000 + i * 3800000)
    c = rect(s, x, Emu(4400000), Emu(3500000), Emu(2000000), CARD1)
    c.line.color.rgb = BORDER
    c.line.width = Pt(1)
    tb(s, Emu(x.emu+300000), Emu(4550000), Emu(2900000), Emu(500000), title, sz=20, b=True, c=accent, a=PP_ALIGN.CENTER)
    tb(s, Emu(x.emu+300000), Emu(5050000), Emu(2900000), Emu(1200000), desc, sz=14, c=LGRAY, a=PP_ALIGN.CENTER)

# SLIDE 4 - Tasks
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Zadachi raboty", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
tasks = [
    ("1.", "Analiz predmetnoj oblasti", "Izuchenie rynka dostavy i trebovanij k avtomatizacii"),
    ("2.", "Proektirovanie BD", "ER-diagramma, 4 tablicy s JSONB-poljami"),
    ("3.", "Razrabotka Frontend", "7 HTML-stranic, adaptivnyj CSS, 4 JS-modulja"),
    ("4.", "Razrabotka Backend", "14 API-endpointov, Express.js, NeDB-obortka"),
    ("5.", "Testirovanie", "API-testy (Mocha, Chai, Supertest), 28 testov"),
    ("6.", "Dokumentirovanie", "9 diagramm dlja diploma, tehnicheskaja dokumentacija"),
]
for i, (num, title, desc) in enumerate(tasks):
    row = i // 3; col = i % 3
    x = Emu(400000 + col * 3900000); y = Emu(1400000 + row * 2600000)
    card = rect(s, x, y, Emu(3600000), Emu(2200000), CARD1)
    card.line.color.rgb = BORDER; card.line.width = Pt(1)
    tb(s, Emu(x.emu+200000), Emu(y.emu+150000), Emu(3200000), Emu(600000), num, sz=36, c=RED)
    tb(s, Emu(x.emu+200000), Emu(y.emu+700000), Emu(3200000), Emu(500000), title, sz=18, b=True, c=DARK)
    tb(s, Emu(x.emu+200000), Emu(y.emu+1200000), Emu(3200000), Emu(800000), desc, sz=13, c=LGRAY)

# SLIDE 5 - Object & Subject
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Obieekt i predmet issledovanija", sz=32, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
obj = rect(s, Emu(400000), Emu(1500000), Emu(5400000), Emu(4600000), RGBColor(0xFD,0xF0,0xF0))
obj.line.color.rgb = RGBColor(0xF5,0xD0,0xD0); obj.line.width = Pt(1)
tb(s, Emu(600000), Emu(1700000), Emu(4900000), Emu(600000), "Obieekt issledovanija", sz=24, b=True, c=RED)
tb(s, Emu(600000), Emu(2400000), Emu(4900000), Emu(1500000),
   "Process prijoma i obrabotki zakazov\nv italjanskom restorane\n\"La Bella Italia\"", sz=18, c=DARK)
tb(s, Emu(1800000), Emu(4000000), Emu(2000000), Emu(1200000), "Pizza", sz=48, a=PP_ALIGN.CENTER)
subj = rect(s, Emu(6200000), Emu(1500000), Emu(5400000), Emu(4600000), RGBColor(0xF0,0xF5,0xFD))
subj.line.color.rgb = RGBColor(0xD0,0xD8,0xF5); subj.line.width = Pt(1)
tb(s, Emu(6400000), Emu(1700000), Emu(4900000), Emu(600000), "Predmet issledovanija", sz=24, b=True, c=RGBColor(0x1A,0x5C,0x8A))
tb(s, Emu(6400000), Emu(2400000), Emu(4900000), Emu(1500000),
   "Veb-prilozhenie dlja avtomatizacii\nprijoma, obrabotki i otslezhivanija\nzakazov s bonusnoj sistemoj", sz=18, c=DARK)
tb(s, Emu(8200000), Emu(4000000), Emu(2000000), Emu(1200000), "Monitor", sz=48, a=PP_ALIGN.CENTER)

# SLIDE 6 - Tech comparison
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Sravnitelnyj analiz tehnologij", sz=32, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
hds = ["Kriterij", "Node.js + Express (nash stek)", "Python + Flask (alternativa)"]
hx = [Emu(400000), Emu(3400000), Emu(7300000)]
hw = [Emu(2800000), Emu(3700000), Emu(4400000)]
for j, (h, px, pw) in enumerate(zip(hds, hx, hw)):
    rect(s, px, Emu(1400000), pw, Emu(600000), DARK)
    tb(s, Emu(px.emu+100000), Emu(1450000), Emu(pw.emu-200000), Emu(500000), h, sz=13, b=True, a=PP_ALIGN.CENTER)
rows = [
    ("Proizvoditelnost", "Vysokaja (asinhronnost)", "Srednjaja"),
    ("Ekologija", "npm (2M+ paketov)", "PyPI (400K+ paketov)"),
    ("Baza dannyh", "PostgreSQL + NeDB", "SQLite"),
    ("Bezopasnost", "Helmet + bcrypt + rate-limit", "Flask-Security"),
    ("Bonusnaja sistema", "5% keshbek, do 30% skidka", "Net vstroennoj"),
    ("Hranenie korziny", "localStorage (bez nagruzki)", "Na servere"),
]
for i, (cr, our, alt) in enumerate(rows):
    ry = Emu(2000000 + i * 700000)
    bgc = RGBColor(0xFA,0xFA,0xFA) if i % 2 == 0 else WHITE
    for j, (txt, px, pw) in enumerate(zip([cr, our, alt], hx, hw)):
        cell = rect(s, px, ry, pw, Emu(650000), bgc)
        cell.line.color.rgb = BORDER; cell.line.width = Pt(0.5)
        cl = DARK if j < 2 else LGRAY
        tb(s, Emu(px.emu+80000), Emu(ry.emu+80000), Emu(pw.emu-160000), Emu(500000), txt, sz=12, b=j==0, c=cl, a=PP_ALIGN.CENTER)

# SLIDE 7 - Competitors
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Obzor sushhestvujushih reshenij", sz=32, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
comps = [
    ("Yandex Eda", ["Bolshaja baza restoranov","Bystraja dostavka","Mobilnoe prilozhenie"],
     ["Vysokaja komissija 30%","Slozhnaja integracija","Net bonusov naprjamuju"]),
    ("Delivery Club", ["Izvestnyj brend","Shirokij ohvat","Raznye sposoby oplaty"],
     ["Vysokaja komissija","Dolgie raschety","Net gibkih skidok"]),
    ("Nashe reshenie La Bella Italia", ["0% komissii","Sobstvennaja BD","Gibkaja bonusnaja sistema","Polnyj kontrol"],
     ["Trebuet hostinga","Nachalnoe napolnenie"]),
]
for i, (name, pros, cons) in enumerate(comps):
    x = Emu(300000 + i * 4000000); y = Emu(1400000)
    is_ours = (i == 2)
    bgc = RGBColor(0xFD,0xF7,0xEC) if is_ours else CARD1
    bdr = GOLD if is_ours else BORDER
    card = rect(s, x, y, Emu(3700000), Emu(5100000), bgc)
    card.line.color.rgb = bdr; card.line.width = Pt(2) if is_ours else Pt(1)
    tb(s, Emu(x.emu+200000), Emu(y.emu+150000), Emu(3300000), Emu(600000), name, sz=16, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, Emu(x.emu+200000), Emu(y.emu+900000), Emu(3300000), Emu(300000), "PLUSY:", sz=12, b=True, c=GREEN)
    for pi, p in enumerate(pros):
        tb(s, Emu(x.emu+200000), Emu(y.emu+1200000+pi*350000), Emu(3300000), Emu(300000), " + %s"%p, sz=11, c=DARK)
    base_cons = Emu(1200000 + len(pros) * 350000 + 100000)
    tb(s, Emu(x.emu+200000), Emu(y.emu+base_cons.emu), Emu(3300000), Emu(300000), "MINUSY:", sz=12, b=True, c=RED)
    for ci, c in enumerate(cons):
        tb(s, Emu(x.emu+200000), Emu(y.emu+base_cons.emu+300000+ci*300000), Emu(3300000), Emu(300000), " - %s"%c, sz=11, c=LGRAY)

# SLIDE 8 - Implementation
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Prakticheskaja realizacija", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
tb(s, Emu(500000), Emu(1400000), Emu(5500000), Emu(500000), "Frontend", sz=24, b=True, c=RED)
lines(s, Emu(500000), Emu(1900000), Emu(5500000), Emu(2000000), [
    {'t': '7 HTML-stranic', 's': 14, 'c': DARK},
    {'t': 'Edinyj CSS-fajl s adaptivnym dizajnom', 's': 14, 'c': DARK},
    {'t': '4 JS-modulja (db, auth, cart, app)', 's': 14, 'c': DARK},
    {'t': 'Korzina v localStorage - bez nagruzki na server', 's': 14, 'c': DARK},
    {'t': 'Maska telefona +7, validacija adresa Moskvy', 's': 14, 'c': DARK},
    {'t': 'Animacii, skeleton-zagruzka', 's': 14, 'c': DARK},
])
tb(s, Emu(6200000), Emu(1400000), Emu(5500000), Emu(500000), "Backend", sz=24, b=True, c=GREEN)
lines(s, Emu(6200000), Emu(1900000), Emu(5500000), Emu(2000000), [
    {'t': 'Node.js + Express.js na portu 3000', 's': 14, 'c': DARK},
    {'t': '14 REST API-endpointov', 's': 14, 'c': DARK},
    {'t': 'NeDB-obortka dlja sovmestimosti', 's': 14, 'c': DARK},
    {'t': 'Helmet (CSP), rate-limit (10 pop/15 min)', 's': 14, 'c': DARK},
    {'t': 'bcrypt dlja hheshirovanija parolej', 's': 14, 'c': DARK},
    {'t': 'express-session dlja sessij', 's': 14, 'c': DARK},
])
tb(s, Emu(500000), Emu(3700000), Emu(11000000), Emu(500000), "Screenshots prilozhenija", sz=20, b=True, c=DARK)
ss = [("Glavnaja", "index.html"),("Menu", "menu.html"),("Korzina", "cart.html"),("Admin", "admin.html")]
for i, (lbl, pg) in enumerate(ss):
    x = Emu(300000 + i * 2950000)
    placeholder(s, x, Emu(4200000), Emu(2700000), Emu(2200000), "[Screenshot]\n%s (%s)" % (lbl, pg))

# SLIDE 9 - Results (diagrams)
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Rezultaty raboty", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
tb(s, Emu(500000), Emu(1300000), Emu(5000000), Emu(400000), "Diagramma arhitektury", sz=18, b=True, c=DARK, a=PP_ALIGN.CENTER)
placeholder(s, Emu(400000), Emu(1700000), Emu(5300000), Emu(3500000),
   "[ARCHITECTURE DIAGRAM]\nVstavte izobrazhenie iz\npapki diploma_diagrams/architecture/")
tb(s, Emu(6200000), Emu(1300000), Emu(5000000), Emu(400000), "ER-diagramma BD", sz=18, b=True, c=DARK, a=PP_ALIGN.CENTER)
placeholder(s, Emu(6100000), Emu(1700000), Emu(5300000), Emu(3500000),
   "[ER DIAGRAM]\nVstavte izobrazhenie iz\npapki diploma_diagrams/er/")
tb(s, Emu(500000), Emu(5500000), Emu(11000000), Emu(500000),
   "Vse 9 diagramm v papke diploma_diagrams/ - gotovy k vstavke v diplom",
   sz=13, c=LGRAY, a=PP_ALIGN.CENTER)

# SLIDE 10 - Testing
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Testirovanie", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
lines(s, Emu(500000), Emu(1400000), Emu(11000000), Emu(800000), [
    {'t': 'API-testirovanie: Mocha + Chai + Supertest - 28 testov', 's': 18, 'c': DARK, 'b': True},
    {'t': 'Menju (4), Registracija (3), Vhod/Vyhod (6), Profil (2), Zakazy (9), Admin (4)', 's': 14, 'c': LGRAY},
])
hds2 = ["Kategorija", "Testov", "Status"]
hx2 = [Emu(500000), Emu(5000000), Emu(8500000)]
hw2 = [Emu(4300000), Emu(3300000), Emu(2700000)]
for j, (h, px, pw) in enumerate(zip(hds2, hx2, hw2)):
    rect(s, px, Emu(2300000), pw, Emu(450000), DARK)
    tb(s, Emu(px.emu+100000), Emu(2360000), Emu(pw.emu-200000), Emu(380000), h, sz=15, b=True, a=PP_ALIGN.CENTER)
trows = [
    ("Publichnoe menu", "4 testa", "Projdeno"),
    ("Registracija", "3 testa", "Projdeno"),
    ("Vhod i vyhod", "6 testov", "Projdeno"),
    ("Profil", "2 testa", "Projdeno"),
    ("Zakazy", "9 testov", "Projdeno"),
    ("Administrirovanie", "4 testa", "Projdeno"),
    ("Statistika", "3 testa", "Projdeno"),
    ("404 i oshibki", "2 testa", "Projdeno"),
]
for i, (cat, cnt, st) in enumerate(trows):
    ry = Emu(2750000 + i * 380000)
    bgc = RGBColor(0xF8,0xFF,0xF0) if i % 2 == 0 else WHITE
    for j, (txt, px, pw) in enumerate(zip([cat, cnt, st], hx2, hw2)):
        cell = rect(s, px, ry, pw, Emu(360000), bgc)
        cell.line.color.rgb = BORDER; cell.line.width = Pt(0.5)
        cl = GREEN if j == 2 else DARK
        bd = True if j > 0 else False
        tb(s, Emu(px.emu+80000), Emu(ry.emu+30000), Emu(pw.emu-160000), Emu(300000), txt, sz=13, b=bd, c=cl, a=PP_ALIGN.CENTER)

# SLIDE 11 - Conclusion
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, 0, 0, SW, Emu(80000), RED)
tb(s, Emu(500000), Emu(300000), Emu(11000000), Emu(900000), "Zakljuchenie", sz=36, b=True, c=DARK)
rect(s, Emu(500000), Emu(1100000), Emu(1500000), Emu(40000), RED)
ach = [
    "Razrabotano polnofunkcionalnoe veb-prilozhenie dlja restorana",
    "Sozdana bonusnaja sistema s keshbekom 5% i skidkoj do 30%",
    "Realizovana panel administratora dlja upravlenija zakazami i menu",
    "Obespechena bezopasnost: Helmet, bcrypt, rate-limit, sessii",
    "Razrabotano 9 UML-diagramm dlja diplomnogo proekta",
    "Provedeno testirovanie: 28 API-testov uspeshno projdeny",
]
for i, txt in enumerate(ach):
    tb(s, Emu(500000), Emu(1300000+i*600000), Emu(11000000), Emu(500000), "+ %s"%txt, sz=17, c=DARK)
card = rect(s, Emu(1500000), Emu(5200000), Emu(9000000), Emu(1200000), RGBColor(0xF5,0xF0,0xE0))
tb(s, Emu(1700000), Emu(5400000), Emu(8600000), Emu(800000),
   "Cel dostignuta: razrabotano veb-prilozhenie, gotovoe k razvertyvaniju\n"
   "i ispolzovaniju v italjanskom restorane \"La Bella Italia\"",
   sz=18, b=True, c=DARK, a=PP_ALIGN.CENTER)

# SLIDE 12 - Thank you
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK)
rect(s, 0, 0, SW, Emu(80000), RED)
rect(s, 0, 0, Emu(200000), SH, RED)
tb(s, Emu(1000000), Emu(1500000), Emu(10000000), Emu(2000000), "Spasibo za vnimanie!", sz=56, b=True, a=PP_ALIGN.CENTER)
tb(s, Emu(1000000), Emu(3200000), Emu(10000000), Emu(1500000), "Pizza Pasta Dolce Vino", sz=48, a=PP_ALIGN.CENTER)
tb(s, Emu(1000000), Emu(4500000), Emu(10000000), Emu(800000),
   "La Bella Italia - Vkus Italii k vashej dveri", sz=24, c=GOLD, a=PP_ALIGN.CENTER)
tb(s, Emu(3000000), Emu(5500000), Emu(6000000), Emu(500000),
   "labellaitalia.ru | +7 (985) 142-82-70 | Moskva, ul. Tverskaja, 12",
   sz=14, c=RGBColor(0x99,0x99,0x99), a=PP_ALIGN.CENTER)

# SAVE
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Prezentacija_La_Bella_Italia.pptx")
prs.save(out)
print("SAVED: " + out)
print("Slides: %d" % len(prs.slides))
sys.exit(0)
