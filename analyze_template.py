from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

prs = Presentation('Maket_prezentatsii_1.pptx')
print(f'Slide width: {prs.slide_width} EMU = {prs.slide_width/914400:.2f} inches')
print(f'Slide height: {prs.slide_height} EMU = {prs.slide_height/914400:.2f} inches')
print(f'Number of slides: {len(prs.slides)}')
print(f'Slide layouts available:')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: {layout.name}')
print()

for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} (layout: {slide.slide_layout.name}) ===')
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        print(f'  Shape: type={stype}, name="{shape.name}"')
        print(f'    pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text
                if text.strip():
                    font_info = ''
                    if para.runs:
                        run = para.runs[0]
                        try:
                            fs = run.font.size
                            fb = run.font.bold
                            fi = run.font.italic
                            fc = 'N/A'
                            try:
                                if run.font.color and run.font.color.rgb:
                                    fc = str(run.font.color.rgb)
                            except:
                                pass
                            font_info = f' [size={fs}, bold={fb}, italic={fi}, color={fc}]'
                        except:
                            pass
                    print(f'    Text: "{text[:100]}"{font_info}')
                    print(f'    Alignment: {para.alignment}')
        if hasattr(shape, 'image'):
            try:
                img = shape.image
                print(f'    Image: {img.content_type}, {len(img.blob)} bytes')
            except:
                pass
    print()
